"""
Document Generation Tools
=========================
Async tools for generating downloadable documents (PPTX, XLSX, DOCX, CSV).

Generated files are saved to ./generated_documents/ and served via an
embedded file server (daemon thread) that auto-starts on import.

The AI includes a ```generated-file``` code block in its response so
the frontend can render a download card.
"""

import os
import json
import asyncio
import socket
import mimetypes
import threading
from typing import Dict, Any, List, Optional
from datetime import datetime
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import unquote, parse_qs, urlparse, urlencode

from langchain_core.tools import tool

# Output directory for generated documents
# Uses env var if set, otherwise defaults to ./generated_documents relative to project root
GENERATED_DOCS_DIR = os.getenv(
    "GENERATED_DOCS_DIR",
    os.path.join(
        os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
        "generated_documents",
    ),
)

# File server port (must match frontend expectations)
FILE_SERVER_PORT = int(os.getenv("FILE_SERVER_PORT", "8090"))

# Skills directory — each sub-folder contains a SKILL.md with YAML frontmatter
SKILLS_DIR = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "skills"
)

# Ensure proper MIME types for Office docs
mimetypes.add_type(
    "application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"
)
mimetypes.add_type(
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"
)
mimetypes.add_type(
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"
)


# =============================================================================
# EMBEDDED FILE SERVER (daemon thread — auto-starts on import)
# =============================================================================

ALLOWED_EXTENSIONS = {".pptx", ".xlsx", ".docx", ".csv", ".pdf", ".txt", ".json", ".md"}

_file_server_started = False
_file_server_lock = threading.Lock()


class _FileRequestHandler(BaseHTTPRequestHandler):
    """Minimal HTTP handler for serving generated document files."""

    def _send_cors_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, PUT, DELETE, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type, Authorization")

    def do_OPTIONS(self):  # noqa: N802
        self.send_response(204)
        self._send_cors_headers()
        self.end_headers()

    def do_GET(self):  # noqa: N802
        path = unquote(self.path)

        # Health check
        if path == "/health":
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self._send_cors_headers()
            self.end_headers()
            self.wfile.write(b'{"status":"ok","service":"embedded-file-server"}')
            return

        # File list
        if path == "/files":
            self._handle_list_files()
            return

        # File download
        if path.startswith("/download/"):
            filename = path[len("/download/"):]
            self._handle_download(filename)
            return

        # Feedback stats
        if path == "/feedback/stats":
            self._handle_feedback_stats()
            return

        # Feedback events
        if path.startswith("/feedback/events"):
            self._handle_feedback_events()
            return

        # Skills list
        if path == "/skills":
            self._handle_list_skills()
            return

        # Single skill detail
        if path.startswith("/skills/") and path.count("/") == 2:
            dir_name = path[len("/skills/"):]
            self._handle_get_skill(dir_name)
            return

        self.send_response(404)
        self.send_header("Content-Type", "application/json")
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(b'{"error":"Not found"}')

    def _handle_list_files(self):
        import glob as _glob

        files = []
        for ext in ALLOWED_EXTENSIONS:
            for fp in _glob.glob(os.path.join(GENERATED_DOCS_DIR, f"*{ext}")):
                stat = os.stat(fp)
                name = os.path.basename(fp)
                files.append({
                    "filename": name,
                    "file_type": os.path.splitext(name)[1].lstrip("."),
                    "size_bytes": stat.st_size,
                    "download_url": f"/download/{name}",
                })
        body = json.dumps({"files": files, "count": len(files)}).encode()
        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(body)

    def _handle_download(self, filename: str):
        # Security: prevent path traversal
        if ".." in filename or "/" in filename or "\\" in filename:
            self.send_response(400)
            self.send_header("Content-Type", "application/json")
            self._send_cors_headers()
            self.end_headers()
            self.wfile.write(b'{"error":"Invalid filename"}')
            return

        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        ext = os.path.splitext(filename)[1].lower()

        if ext not in ALLOWED_EXTENSIONS:
            self.send_response(403)
            self.send_header("Content-Type", "application/json")
            self._send_cors_headers()
            self.end_headers()
            self.wfile.write(b'{"error":"File type not allowed"}')
            return

        if not os.path.isfile(filepath):
            self.send_response(404)
            self.send_header("Content-Type", "application/json")
            self._send_cors_headers()
            self.end_headers()
            self.wfile.write(b'{"error":"File not found"}')
            return

        content_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"

        with open(filepath, "rb") as f:
            data = f.read()

        self.send_response(200)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
        self.send_header("Content-Length", str(len(data)))
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(data)

    # ----- POST handler -----

    def do_POST(self):  # noqa: N802
        path = unquote(self.path)

        if path == "/feedback":
            self._handle_feedback_post()
            return

        if path == "/notify/email":
            self._handle_notify_email()
            return

        if path == "/notify/test-email":
            self._handle_notify_test_email()
            return

        # Create a new skill
        if path == "/skills":
            self._handle_create_skill()
            return

        # Upload a skill .md file
        if path == "/skills/upload":
            self._handle_upload_skill()
            return

        # Reload skills into the running agent
        if path == "/skills/reload":
            self._handle_reload_skills()
            return

        # CORS preflight is handled by do_OPTIONS; any other POST is 404
        self.send_response(404)
        self.send_header("Content-Type", "application/json")
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(b'{"error":"Not found"}')

    # ----- PUT handler -----

    def do_PUT(self):  # noqa: N802
        path = unquote(self.path)

        # Update an existing skill
        if path.startswith("/skills/") and path.count("/") == 2:
            dir_name = path[len("/skills/"):]
            self._handle_update_skill(dir_name)
            return

        self.send_response(404)
        self.send_header("Content-Type", "application/json")
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(b'{"error":"Not found"}')

    # ----- DELETE handler -----

    def do_DELETE(self):  # noqa: N802
        path = unquote(self.path)

        # Delete a skill
        if path.startswith("/skills/") and path.count("/") == 2:
            dir_name = path[len("/skills/"):]
            self._handle_delete_skill(dir_name)
            return

        self.send_response(404)
        self.send_header("Content-Type", "application/json")
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(b'{"error":"Not found"}')

    # ----- Feedback handlers -----

    def _handle_feedback_post(self):
        """Record a user feedback event (POST /feedback)."""
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length) if content_length > 0 else b"{}"
            data = json.loads(body)
        except Exception:
            self._send_json(400, {"error": "Invalid JSON body"})
            return

        signal_str = data.get("signal", "")
        thread_id = data.get("thread_id", "unknown")
        context = data.get("context", "")
        domain = data.get("domain", "")
        metadata = data.get("metadata", {})

        try:
            from sdtm_pipeline.deepagents.feedback import FeedbackSignal, get_feedback_collector

            try:
                signal = FeedbackSignal(signal_str)
            except ValueError:
                valid = [s.value for s in FeedbackSignal]
                self._send_json(400, {"error": f"Invalid signal: {signal_str}", "valid_signals": valid})
                return

            collector = get_feedback_collector()
            event = collector.record(
                signal=signal,
                thread_id=thread_id,
                user_query=context,
                domain=domain.upper() if domain else None,
                metadata=metadata,
            )

            self._send_json(200, {
                "success": True,
                "event_id": event.event_id,
                "signal": signal_str,
                "sentiment": event.sentiment.value,
            })
        except ImportError:
            self._send_json(503, {"error": "Feedback module not available"})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_feedback_stats(self):
        """Return aggregated feedback analytics (GET /feedback/stats)."""
        try:
            from sdtm_pipeline.deepagents.learning_store import get_learning_store

            store = get_learning_store()
            metrics = store.compute_metrics()
            self._send_json(200, metrics)
        except ImportError:
            self._send_json(503, {"error": "Learning store not available"})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_feedback_events(self):
        """Return recent feedback events (GET /feedback/events?limit=50&thread_id=xxx)."""
        try:
            from sdtm_pipeline.deepagents.learning_store import get_learning_store

            store = get_learning_store()

            # Parse query parameters
            parsed = urlparse(self.path)
            params = parse_qs(parsed.query)
            limit = int(params.get("limit", ["50"])[0])
            thread_id = params.get("thread_id", [None])[0]

            events = store.read_events(limit=limit, thread_id=thread_id)
            self._send_json(200, {
                "events": [e.to_dict() for e in events],
                "count": len(events),
            })
        except ImportError:
            self._send_json(503, {"error": "Learning store not available"})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    # ----- Skills CRUD handlers -----

    def _parse_skill_frontmatter(self, content: str):
        """Parse YAML frontmatter from a SKILL.md file. Returns (name, description, body)."""
        name, description, body = "", "", content
        if content.startswith("---"):
            parts = content.split("---", 2)
            if len(parts) >= 3:
                frontmatter = parts[1].strip()
                body = parts[2].strip()
                for line in frontmatter.splitlines():
                    if line.startswith("name:"):
                        name = line[len("name:"):].strip()
                    elif line.startswith("description:"):
                        description = line[len("description:"):].strip()
        return name, description, body

    def _build_skill_md(self, name: str, description: str, body: str) -> str:
        """Build a SKILL.md file with YAML frontmatter."""
        return f"---\nname: {name}\ndescription: {description}\n---\n\n{body}\n"

    def _slugify(self, text: str) -> str:
        """Convert text to a directory-safe slug."""
        import re
        slug = text.lower().strip()
        slug = re.sub(r'[^a-z0-9\s-]', '', slug)
        slug = re.sub(r'[\s_]+', '-', slug)
        slug = re.sub(r'-+', '-', slug)
        return slug.strip('-')

    def _reload_agent_skills(self):
        """Reload skills into the running LangGraph agent."""
        try:
            from sdtm_pipeline.langgraph_chat.graph import reload_skills
            reload_skills()
        except Exception as e:
            print(f"[Skills] Warning: Could not reload skills: {e}")

    def _handle_list_skills(self):
        """List all skills (GET /skills)."""
        try:
            skills = []
            if os.path.isdir(SKILLS_DIR):
                for dir_name in sorted(os.listdir(SKILLS_DIR)):
                    skill_dir = os.path.join(SKILLS_DIR, dir_name)
                    skill_file = os.path.join(skill_dir, "SKILL.md")
                    if os.path.isdir(skill_dir) and os.path.isfile(skill_file):
                        with open(skill_file, "r", encoding="utf-8") as f:
                            content = f.read()
                        name, description, _ = self._parse_skill_frontmatter(content)
                        skills.append({
                            "name": name or dir_name,
                            "description": description,
                            "dirName": dir_name,
                        })
            self._send_json(200, {"skills": skills, "count": len(skills)})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_get_skill(self, dir_name: str):
        """Get full skill detail (GET /skills/<dir_name>)."""
        skill_file = os.path.join(SKILLS_DIR, dir_name, "SKILL.md")
        if not os.path.isfile(skill_file):
            self._send_json(404, {"error": f"Skill '{dir_name}' not found"})
            return
        try:
            with open(skill_file, "r", encoding="utf-8") as f:
                raw = f.read()
            name, description, body = self._parse_skill_frontmatter(raw)
            self._send_json(200, {
                "name": name or dir_name,
                "description": description,
                "dirName": dir_name,
                "content": body,
            })
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_create_skill(self):
        """Create a new skill (POST /skills)."""
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length) if content_length > 0 else b"{}"
            data = json.loads(body)
        except Exception:
            self._send_json(400, {"error": "Invalid JSON body"})
            return

        name = data.get("name", "").strip()
        description = data.get("description", "").strip()
        content = data.get("content", "").strip()

        if not name:
            self._send_json(400, {"error": "Missing 'name'"})
            return

        dir_name = self._slugify(name)
        if not dir_name:
            self._send_json(400, {"error": "Invalid skill name"})
            return

        skill_dir = os.path.join(SKILLS_DIR, dir_name)
        if os.path.exists(skill_dir):
            self._send_json(409, {"error": f"Skill '{dir_name}' already exists"})
            return

        try:
            os.makedirs(skill_dir, exist_ok=True)
            skill_file = os.path.join(skill_dir, "SKILL.md")
            with open(skill_file, "w", encoding="utf-8") as f:
                f.write(self._build_skill_md(name, description, content))
            self._reload_agent_skills()
            self._send_json(201, {"success": True, "dirName": dir_name})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_upload_skill(self):
        """Upload a .md file as a new skill (POST /skills/upload)."""
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            raw_body = self.rfile.read(content_length) if content_length > 0 else b""
        except Exception:
            self._send_json(400, {"error": "Could not read request body"})
            return

        content_type = self.headers.get("Content-Type", "")

        # Handle JSON upload (base64 or raw text)
        if "application/json" in content_type:
            try:
                data = json.loads(raw_body)
                file_content = data.get("content", "")
                file_name = data.get("fileName", "uploaded-skill")
            except Exception:
                self._send_json(400, {"error": "Invalid JSON body"})
                return
        else:
            # Treat as raw text/markdown
            file_content = raw_body.decode("utf-8", errors="replace")
            file_name = "uploaded-skill"

        if not file_content.strip():
            self._send_json(400, {"error": "Empty file content"})
            return

        # Parse frontmatter if present
        name, description, body = self._parse_skill_frontmatter(file_content)
        if not name:
            # Use filename (minus extension) as name
            name = os.path.splitext(file_name)[0]
        if not body:
            body = file_content

        dir_name = self._slugify(name)
        if not dir_name:
            dir_name = "uploaded-skill"

        # Handle conflicts by appending a suffix
        base_dir_name = dir_name
        counter = 1
        while os.path.exists(os.path.join(SKILLS_DIR, dir_name)):
            dir_name = f"{base_dir_name}-{counter}"
            counter += 1

        try:
            skill_dir = os.path.join(SKILLS_DIR, dir_name)
            os.makedirs(skill_dir, exist_ok=True)
            skill_file = os.path.join(skill_dir, "SKILL.md")
            with open(skill_file, "w", encoding="utf-8") as f:
                f.write(self._build_skill_md(name, description, body))
            self._reload_agent_skills()
            self._send_json(201, {"success": True, "dirName": dir_name, "name": name})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_reload_skills(self):
        """Manually trigger a skills reload (POST /skills/reload)."""
        try:
            self._reload_agent_skills()
            self._send_json(200, {"success": True})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_update_skill(self, dir_name: str):
        """Update an existing skill (PUT /skills/<dir_name>)."""
        skill_dir = os.path.join(SKILLS_DIR, dir_name)
        skill_file = os.path.join(skill_dir, "SKILL.md")
        if not os.path.isfile(skill_file):
            self._send_json(404, {"error": f"Skill '{dir_name}' not found"})
            return

        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length) if content_length > 0 else b"{}"
            data = json.loads(body)
        except Exception:
            self._send_json(400, {"error": "Invalid JSON body"})
            return

        name = data.get("name", "").strip()
        description = data.get("description", "").strip()
        content = data.get("content", "").strip()

        if not name:
            self._send_json(400, {"error": "Missing 'name'"})
            return

        try:
            with open(skill_file, "w", encoding="utf-8") as f:
                f.write(self._build_skill_md(name, description, content))
            self._reload_agent_skills()
            self._send_json(200, {"success": True, "dirName": dir_name})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    def _handle_delete_skill(self, dir_name: str):
        """Delete a skill directory (DELETE /skills/<dir_name>)."""
        import shutil

        skill_dir = os.path.join(SKILLS_DIR, dir_name)
        if not os.path.isdir(skill_dir):
            self._send_json(404, {"error": f"Skill '{dir_name}' not found"})
            return

        try:
            shutil.rmtree(skill_dir)
            self._reload_agent_skills()
            self._send_json(200, {"success": True})
        except Exception as e:
            self._send_json(500, {"error": str(e)})

    # ----- Email notification handlers -----

    def _get_gmail_oauth_config(self):
        """Read Gmail OAuth2 settings from environment variables."""
        client_id = os.getenv("GMAIL_CLIENT_ID")
        client_secret = os.getenv("GMAIL_CLIENT_SECRET")
        refresh_token = os.getenv("GMAIL_REFRESH_TOKEN")
        from_addr = os.getenv("SMTP_FROM", "")
        if not all([client_id, client_secret, refresh_token, from_addr]):
            return None
        return {
            "client_id": client_id,
            "client_secret": client_secret,
            "refresh_token": refresh_token,
            "from": from_addr,
        }

    def _get_smtp_config(self):
        """Read basic SMTP settings from environment variables (fallback)."""
        host = os.getenv("SMTP_HOST")
        port = int(os.getenv("SMTP_PORT", "587"))
        user = os.getenv("SMTP_USER")
        password = os.getenv("SMTP_PASSWORD")
        from_addr = os.getenv("SMTP_FROM", user or "")
        if not host or not user or not password:
            return None
        return {"host": host, "port": port, "user": user, "password": password, "from": from_addr}

    def _refresh_oauth2_access_token(self, client_id: str, client_secret: str, refresh_token: str) -> str:
        """Exchange a Gmail OAuth2 refresh token for a fresh access token."""
        import urllib.request as _req

        data = urlencode({
            "client_id": client_id,
            "client_secret": client_secret,
            "refresh_token": refresh_token,
            "grant_type": "refresh_token",
        }).encode()

        req = _req.Request(
            "https://oauth2.googleapis.com/token",
            data=data,
            method="POST",
        )
        req.add_header("Content-Type", "application/x-www-form-urlencoded")

        with _req.urlopen(req, timeout=15) as resp:
            body = json.loads(resp.read().decode())

        access_token = body.get("access_token")
        if not access_token:
            raise RuntimeError(f"OAuth2 token refresh failed: {body}")
        return access_token

    def _send_email(self, to: str, subject: str, html_body: str):
        """Send an email. Tries Gmail OAuth2 first, falls back to basic SMTP."""
        import smtplib
        import base64
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart

        # --- Try Gmail OAuth2 (XOAUTH2) ---
        oauth_cfg = self._get_gmail_oauth_config()
        if oauth_cfg:
            access_token = self._refresh_oauth2_access_token(
                oauth_cfg["client_id"],
                oauth_cfg["client_secret"],
                oauth_cfg["refresh_token"],
            )

            msg = MIMEMultipart("alternative")
            msg["From"] = oauth_cfg["from"]
            msg["To"] = to
            msg["Subject"] = subject
            msg.attach(MIMEText(html_body, "html"))

            # Build XOAUTH2 string: "user=<email>\x01auth=Bearer <token>\x01\x01"
            auth_string = f"user={oauth_cfg['from']}\x01auth=Bearer {access_token}\x01\x01"
            auth_b64 = base64.b64encode(auth_string.encode()).decode()

            with smtplib.SMTP("smtp.gmail.com", 587, timeout=15) as server:
                server.ehlo()
                server.starttls()
                server.ehlo()
                code, resp_msg = server.docmd("AUTH", f"XOAUTH2 {auth_b64}")
                if code not in (235, 250):
                    raise RuntimeError(
                        f"XOAUTH2 authentication failed ({code}): {resp_msg.decode()}"
                    )
                server.sendmail(oauth_cfg["from"], [to], msg.as_string())
            return

        # --- Fallback: basic SMTP with password ---
        smtp_cfg = self._get_smtp_config()
        if smtp_cfg:
            msg = MIMEMultipart("alternative")
            msg["From"] = smtp_cfg["from"]
            msg["To"] = to
            msg["Subject"] = subject
            msg.attach(MIMEText(html_body, "html"))

            with smtplib.SMTP(smtp_cfg["host"], smtp_cfg["port"], timeout=15) as server:
                server.ehlo()
                server.starttls()
                server.ehlo()
                server.login(smtp_cfg["user"], smtp_cfg["password"])
                server.sendmail(smtp_cfg["from"], [to], msg.as_string())
            return

        raise EnvironmentError(
            "Email not configured. Set GMAIL_CLIENT_ID, GMAIL_CLIENT_SECRET, "
            "GMAIL_REFRESH_TOKEN, and SMTP_FROM environment variables. "
            "Run `python setup_gmail_oauth.py` to obtain a refresh token."
        )

    def _handle_notify_email(self):
        """Send a task-completion notification email (POST /notify/email)."""
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length) if content_length > 0 else b"{}"
            data = json.loads(body)
        except Exception:
            self._send_json(400, {"error": "Invalid JSON body"})
            return

        to = data.get("to", "")
        task_name = data.get("task_name", "Unknown Task")
        output = data.get("output", "")

        if not to:
            self._send_json(400, {"error": "Missing 'to' email address"})
            return

        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        subject = f"Task Completed: {task_name}"
        html = (
            f"<div style='font-family:sans-serif;max-width:600px;margin:0 auto;'>"
            f"<h2 style='color:#1a1a1a;'>Scheduled Task Completed</h2>"
            f"<p><strong>Task:</strong> {task_name}</p>"
            f"<p><strong>Completed at:</strong> {now}</p>"
            f"<hr style='border:none;border-top:1px solid #e5e5e5;margin:16px 0;'/>"
            f"<h3 style='color:#333;'>Output Summary</h3>"
            f"<div style='background:#f9f9f9;padding:12px 16px;border-radius:8px;"
            f"border:1px solid #e5e5e5;white-space:pre-wrap;font-size:14px;'>"
            f"{output[:2000]}"
            f"</div>"
            f"<p style='color:#888;font-size:12px;margin-top:24px;'>"
            f"Sent by SDTM Agent &middot; Automated notification</p>"
            f"</div>"
        )

        try:
            self._send_email(to, subject, html)
            self._send_json(200, {"success": True})
        except EnvironmentError as e:
            self._send_json(503, {"success": False, "error": str(e)})
        except Exception as e:
            self._send_json(500, {"success": False, "error": str(e)})

    def _handle_notify_test_email(self):
        """Send a test notification email (POST /notify/test-email)."""
        try:
            content_length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(content_length) if content_length > 0 else b"{}"
            data = json.loads(body)
        except Exception:
            self._send_json(400, {"error": "Invalid JSON body"})
            return

        to = data.get("to", "")
        if not to:
            self._send_json(400, {"error": "Missing 'to' email address"})
            return

        subject = "Test Notification from SDTM Agent"
        html = (
            "<div style='font-family:sans-serif;max-width:600px;margin:0 auto;'>"
            "<h2 style='color:#1a1a1a;'>Test Notification</h2>"
            "<p>This is a test email from your SDTM Agent. "
            "If you received this, email notifications are working correctly.</p>"
            "<p style='color:#888;font-size:12px;margin-top:24px;'>"
            "Sent by SDTM Agent &middot; Test notification</p>"
            "</div>"
        )

        try:
            self._send_email(to, subject, html)
            self._send_json(200, {"success": True})
        except EnvironmentError as e:
            self._send_json(503, {"success": False, "error": str(e)})
        except Exception as e:
            self._send_json(500, {"success": False, "error": str(e)})

    # ----- Utility -----

    def _send_json(self, status_code: int, data: dict):
        """Send a JSON response with CORS headers."""
        body = json.dumps(data).encode()
        self.send_response(status_code)
        self.send_header("Content-Type", "application/json")
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(body)

    def log_message(self, format, *args):
        """Suppress default request logging to keep LangGraph output clean."""
        pass


def _is_port_in_use(port: int) -> bool:
    """Check if a port is already in use."""
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        return s.connect_ex(("127.0.0.1", port)) == 0


def _start_file_server():
    """Start the embedded file server in a daemon thread (idempotent)."""
    global _file_server_started

    with _file_server_lock:
        if _file_server_started:
            return

        port = FILE_SERVER_PORT

        # Don't start if something else is already listening (e.g. standalone file_server.py)
        if _is_port_in_use(port):
            print(f"[Document Tools] File server port {port} already in use — skipping embedded server")
            _file_server_started = True
            return

        try:
            os.makedirs(GENERATED_DOCS_DIR, exist_ok=True)
            server = HTTPServer(("0.0.0.0", port), _FileRequestHandler)
            thread = threading.Thread(target=server.serve_forever, daemon=True)
            thread.start()
            print(f"[Document Tools] Embedded file server started on port {port}")
            print(f"[Document Tools] Serving files from: {GENERATED_DOCS_DIR}")
            _file_server_started = True
        except Exception as e:
            print(f"[Document Tools] Warning: Could not start file server: {e}")


# Auto-start the file server when this module is imported by LangGraph
_start_file_server()


# =============================================================================
# HELPERS
# =============================================================================

def _get_download_url(filename: str) -> str:
    """Build the download URL for a generated file."""
    return f"/download/{filename}"


def _ensure_output_dir():
    """Create the output directory if it doesn't exist."""
    os.makedirs(GENERATED_DOCS_DIR, exist_ok=True)


def _sanitize_filename(name: str) -> str:
    """Sanitize a filename, keeping only safe characters."""
    safe = "".join(
        c if c.isalnum() or c in ("-", "_", ".") else "_" for c in name
    )
    # Remove consecutive underscores
    while "__" in safe:
        safe = safe.replace("__", "_")
    return safe.strip("_")


def _unique_filename(base_name: str, ext: str) -> str:
    """Generate a unique filename with timestamp to avoid collisions."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    sanitized = _sanitize_filename(base_name)
    return f"{sanitized}_{timestamp}.{ext}"


# =============================================================================
# PRESENTATION (PPTX) TOOL
# =============================================================================

@tool
async def generate_presentation(
    title: str,
    slides: List[Dict[str, str]],
    subtitle: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Generate a PowerPoint presentation (.pptx) with the given slides.

    Args:
        title: Presentation title (used for the title slide and filename).
        slides: List of slide dicts. Each slide has:
            - title (str): Slide title/heading
            - content (str): Slide body text (supports bullet points with newlines)
            - layout (str, optional): "title", "content", "two_column", or "blank"
        subtitle: Optional subtitle for the title slide.

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    def _build_pptx():
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # --- Colour palette ---
        PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)     # deep navy
        ACCENT = RGBColor(0x00, 0x7A, 0xFF)       # vibrant blue
        LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)     # light grey-blue
        TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
        TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
        SUBTLE = RGBColor(0x6B, 0x72, 0x80)

        def _set_slide_bg(slide, color):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _add_shape_fill(slide, left, top, width, height, color):
            from pptx.util import Emu as _E
            shape = slide.shapes.add_shape(
                1, left, top, width, height  # MSO_SHAPE.RECTANGLE = 1
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        # ── Title Slide ──
        slide_layout = prs.slide_layouts[6]  # blank
        sl = prs.slides.add_slide(slide_layout)
        _set_slide_bg(sl, PRIMARY)

        # Accent bar
        _add_shape_fill(sl, Inches(0.8), Inches(2.8), Inches(1.5), Inches(0.06), ACCENT)

        # Title
        txBox = sl.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(10), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = SUBTLE
            p2.space_before = Pt(12)

        # ── Content Slides ──
        for idx, slide_data in enumerate(slides):
            sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            _set_slide_bg(sl, LIGHT_BG)

            # Top accent bar
            _add_shape_fill(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

            # Slide number
            txNum = sl.shapes.add_textbox(
                prs.slide_width - Inches(1), prs.slide_height - Inches(0.5),
                Inches(0.7), Inches(0.3),
            )
            pNum = txNum.text_frame.paragraphs[0]
            pNum.text = f"{idx + 1}"
            pNum.font.size = Pt(11)
            pNum.font.color.rgb = SUBTLE
            pNum.alignment = PP_ALIGN.RIGHT

            # Slide title
            txTitle = sl.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
            tf = txTitle.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", f"Slide {idx + 1}")
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = TEXT_DARK

            # Divider line
            _add_shape_fill(sl, Inches(0.8), Inches(1.35), Inches(1.2), Inches(0.04), ACCENT)

            # Content body
            content = slide_data.get("content", "")
            layout = slide_data.get("layout", "content")

            if layout == "two_column" and "\n---\n" in content:
                col1, col2 = content.split("\n---\n", 1)
                for ci, col_text in enumerate([col1, col2]):
                    left = Inches(0.8) if ci == 0 else Inches(6.8)
                    txBox = sl.shapes.add_textbox(left, Inches(1.7), Inches(5.2), Inches(5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for line in col_text.strip().split("\n"):
                        line = line.strip()
                        if not line:
                            continue
                        pg = tf.add_paragraph()
                        if line.startswith("- ") or line.startswith("* "):
                            pg.text = line[2:]
                            pg.level = 0
                        else:
                            pg.text = line
                        pg.font.size = Pt(16)
                        pg.font.color.rgb = TEXT_DARK
                        pg.space_after = Pt(6)
            else:
                txBox = sl.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for line in content.strip().split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    pg = tf.add_paragraph()
                    if line.startswith("- ") or line.startswith("* "):
                        pg.text = f"\u2022  {line[2:]}"
                    elif line.startswith("## "):
                        pg.text = line[3:]
                        pg.font.bold = True
                        pg.font.size = Pt(20)
                        pg.font.color.rgb = ACCENT
                        pg.space_before = Pt(12)
                        continue
                    else:
                        pg.text = line
                    pg.font.size = Pt(16)
                    pg.font.color.rgb = TEXT_DARK
                    pg.space_after = Pt(6)

        # Save
        _ensure_output_dir()
        filename = _unique_filename(title, "pptx")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        prs.save(filepath)
        return filepath, filename

    filepath, filename = await asyncio.to_thread(_build_pptx)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "pptx",
        "size_bytes": size_bytes,
        "description": f"PowerPoint presentation with {len(slides)} slides: {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# EXCEL (XLSX) TOOL
# =============================================================================

@tool
async def generate_excel(
    title: str,
    sheets: List[Dict[str, Any]],
) -> Dict[str, Any]:
    """
    Generate an Excel workbook (.xlsx) with one or more sheets.

    Args:
        title: Workbook title (used for the filename).
        sheets: List of sheet dicts. Each sheet has:
            - name (str): Sheet/tab name
            - headers (list[str]): Column headers
            - rows (list[list]): Row data (list of lists)

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    def _build_xlsx():
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = Workbook()
        # Remove default sheet
        if wb.active:
            wb.remove(wb.active)

        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill(start_color="007AFF", end_color="007AFF", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell_alignment = Alignment(vertical="top", wrap_text=True)
        thin_border = Border(
            left=Side(style="thin", color="D0D5DD"),
            right=Side(style="thin", color="D0D5DD"),
            top=Side(style="thin", color="D0D5DD"),
            bottom=Side(style="thin", color="D0D5DD"),
        )
        alt_fill = PatternFill(start_color="F0F4F8", end_color="F0F4F8", fill_type="solid")

        total_rows = 0
        for sheet_data in sheets:
            sheet_name = sheet_data.get("name", "Sheet")[:31]  # Excel limit
            ws = wb.create_sheet(title=sheet_name)

            headers = sheet_data.get("headers", [])
            rows = sheet_data.get("rows", [])

            # Write headers
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
                cell.border = thin_border

            # Write data rows
            for row_idx, row_data in enumerate(rows, 2):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.alignment = cell_alignment
                    cell.border = thin_border
                    if row_idx % 2 == 0:
                        cell.fill = alt_fill

            total_rows += len(rows)

            # Auto-size columns (approximate)
            for col_idx, header in enumerate(headers, 1):
                max_len = len(str(header))
                for row_data in rows:
                    if col_idx - 1 < len(row_data):
                        max_len = max(max_len, len(str(row_data[col_idx - 1])))
                ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = min(max_len + 4, 50)

            # Freeze header row
            ws.freeze_panes = "A2"

        _ensure_output_dir()
        filename = _unique_filename(title, "xlsx")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        wb.save(filepath)
        return filepath, filename, total_rows, len(sheets)

    filepath, filename, total_rows, num_sheets = await asyncio.to_thread(_build_xlsx)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "xlsx",
        "size_bytes": size_bytes,
        "description": f"Excel workbook with {num_sheets} sheet(s), {total_rows} data rows: {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# WORD DOCUMENT (DOCX) TOOL
# =============================================================================

@tool
async def generate_word_document(
    title: str,
    sections: List[Dict[str, Any]],
) -> Dict[str, Any]:
    """
    Generate a Word document (.docx) with structured sections.

    Args:
        title: Document title (appears as heading and filename).
        sections: List of section dicts. Each section has:
            - heading (str): Section heading text
            - content (str): Section body (plain text or markdown-like bullet points)
            - level (int, optional): Heading level 1-3, defaults to 1

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    def _build_docx():
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Styles
        style = doc.styles["Normal"]
        font = style.font
        font.name = "Calibri"
        font.size = Pt(11)
        font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # Title
        title_para = doc.add_heading(title, level=0)
        for run in title_para.runs:
            run.font.color.rgb = RGBColor(0x00, 0x7A, 0xFF)

        # Date line
        date_para = doc.add_paragraph(
            f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        )
        date_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        for run in date_para.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
            run.font.italic = True

        doc.add_paragraph("")  # spacer

        # Sections
        for section in sections:
            heading = section.get("heading", "Section")
            content = section.get("content", "")
            level = section.get("level", 1)

            doc.add_heading(heading, level=min(level, 4))

            # Parse content - support bullet points
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    continue
                if line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style="List Bullet")
                elif line.startswith("1.") or line.startswith("2.") or line.startswith("3."):
                    # Strip the number prefix
                    text = line.split(".", 1)[1].strip() if "." in line else line
                    doc.add_paragraph(text, style="List Number")
                else:
                    doc.add_paragraph(line)

        _ensure_output_dir()
        filename = _unique_filename(title, "docx")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        doc.save(filepath)
        return filepath, filename

    filepath, filename = await asyncio.to_thread(_build_docx)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "docx",
        "size_bytes": size_bytes,
        "description": f"Word document with {len(sections)} section(s): {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# CSV FILE TOOL
# =============================================================================

@tool
async def generate_csv_file(
    title: str,
    headers: List[str],
    rows: List[List[Any]],
) -> Dict[str, Any]:
    """
    Generate a CSV file from tabular data.

    Args:
        title: File title (used for the filename).
        headers: List of column header strings.
        rows: List of row data (each row is a list of values).

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    import pandas as pd

    def _build_csv():
        import pandas as pd

        df = pd.DataFrame(rows, columns=headers)
        _ensure_output_dir()
        filename = _unique_filename(title, "csv")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        df.to_csv(filepath, index=False)
        return filepath, filename, len(df)

    filepath, filename, row_count = await asyncio.to_thread(_build_csv)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "csv",
        "size_bytes": size_bytes,
        "description": f"CSV file with {row_count} rows, {len(headers)} columns: {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# PDF DOCUMENT TOOL
# =============================================================================

@tool
async def generate_pdf(
    title: str,
    sections: List[Dict[str, Any]],
) -> Dict[str, Any]:
    """
    Generate a PDF document with structured sections.

    Args:
        title: Document title (appears as heading and filename).
        sections: List of section dicts. Each section has:
            - heading (str): Section heading text
            - content (str): Section body (plain text or bullet points with newlines)
            - level (int, optional): Heading level 1-3, defaults to 1

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    def _build_pdf():
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # Title
        pdf.set_font("Helvetica", "B", 24)
        pdf.set_text_color(0, 122, 255)
        pdf.cell(0, 14, title, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)

        # Date
        pdf.set_font("Helvetica", "I", 10)
        pdf.set_text_color(107, 114, 128)
        pdf.cell(0, 6, f"Generated: {datetime.now().strftime('%B %d, %Y')}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(6)

        # Sections
        for section in sections:
            heading = section.get("heading", "Section")
            content = section.get("content", "")
            level = section.get("level", 1)

            # Heading
            size = {1: 16, 2: 14, 3: 12}.get(level, 14)
            pdf.set_font("Helvetica", "B", size)
            pdf.set_text_color(26, 26, 46)
            pdf.cell(0, 10, heading, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

            # Content
            pdf.set_font("Helvetica", "", 11)
            pdf.set_text_color(51, 51, 51)
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    pdf.ln(3)
                    continue
                if line.startswith("- ") or line.startswith("* "):
                    pdf.cell(8)  # indent
                    pdf.multi_cell(0, 6, f"\u2022  {line[2:]}")
                else:
                    pdf.multi_cell(0, 6, line)
            pdf.ln(4)

        _ensure_output_dir()
        filename = _unique_filename(title, "pdf")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        pdf.output(filepath)
        return filepath, filename

    filepath, filename = await asyncio.to_thread(_build_pdf)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "pdf",
        "size_bytes": size_bytes,
        "description": f"PDF document with {len(sections)} section(s): {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# MARKDOWN FILE TOOL
# =============================================================================

@tool
async def generate_markdown_file(
    title: str,
    content: str,
) -> Dict[str, Any]:
    """
    Generate a Markdown (.md) file.

    Args:
        title: Document title (used for filename and top heading).
        content: Full markdown content. You can use headings, bullet points,
                 tables, code blocks, etc. The title is prepended as a # heading.

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    def _build_md():
        full_content = f"# {title}\n\n_Generated: {datetime.now().strftime('%B %d, %Y')}_\n\n{content}"
        _ensure_output_dir()
        filename = _unique_filename(title, "md")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(full_content)
        return filepath, filename, len(content.split("\n"))

    filepath, filename, line_count = await asyncio.to_thread(_build_md)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "md",
        "size_bytes": size_bytes,
        "description": f"Markdown document ({line_count} lines): {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# PLAIN TEXT FILE TOOL
# =============================================================================

@tool
async def generate_text_file(
    title: str,
    content: str,
) -> Dict[str, Any]:
    """
    Generate a plain text (.txt) file.

    Args:
        title: Document title (used for filename and header line).
        content: Full text content to write to the file.

    Returns:
        Metadata dict with filename, file_type, size_bytes, description, download_url.

    IMPORTANT: After calling this tool, include the returned metadata in your
    response using a ```generated-file``` code block so the user can download it.
    """
    def _build_txt():
        header = f"{title}\n{'=' * len(title)}\nGenerated: {datetime.now().strftime('%B %d, %Y')}\n\n"
        full_content = header + content
        _ensure_output_dir()
        filename = _unique_filename(title, "txt")
        filepath = os.path.join(GENERATED_DOCS_DIR, filename)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(full_content)
        return filepath, filename, len(content.split("\n"))

    filepath, filename, line_count = await asyncio.to_thread(_build_txt)
    size_bytes = os.path.getsize(filepath)

    return {
        "success": True,
        "filename": filename,
        "file_type": "txt",
        "size_bytes": size_bytes,
        "description": f"Text file ({line_count} lines): {title}",
        "download_url": _get_download_url(filename),
        "instruction": (
            "Include the following generated-file code block in your response "
            "so the user can download the file."
        ),
    }


# =============================================================================
# EXPORTS
# =============================================================================

DOCUMENT_TOOLS = [
    generate_presentation,
    generate_excel,
    generate_word_document,
    generate_csv_file,
    generate_pdf,
    generate_markdown_file,
    generate_text_file,
]
